import os
import re

from pptx import Presentation

DATA_DIR = os.path.join(os.path.dirname(__file__), "data")


def get_ppt_text(filepath):
    prs = Presentation(filepath)
    parts = []
    for slide in prs.slides:
        for shape in slide.shapes:
            if shape.has_table:
                for row in shape.table.rows:
                    for cell in row.cells:
                        parts.append(cell.text)
            elif shape.has_text_frame:
                parts.append(shape.text_frame.text)
    text = " ".join(parts)
    text = re.sub(r"\s+", " ", text)
    return text.strip()


def save_ppt_text(text, filename):
    os.makedirs(DATA_DIR, exist_ok=True)
    if not filename.endswith(".txt"):
        filename = filename + ".txt"
    path = os.path.join(DATA_DIR, filename)
    with open(path, "w", encoding="utf-8") as f:
        f.write(text)


def load_ppt_text(filename):
    if not filename.endswith(".txt"):
        filename = filename + ".txt"
    path = os.path.join(DATA_DIR, filename)
    with open(path, encoding="utf-8") as f:
        return f.read()
